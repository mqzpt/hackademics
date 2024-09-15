from flask import Flask, jsonify, request
from flask_cors import CORS
import cohere
import pdfplumber
import docx
import assemblyai as aai
from pptx import Presentation
import os
import random

app = Flask(__name__)
CORS(app)

co = cohere.Client('s6QrVVNoTa8VaUiBUhWvVQk1jU0qu2E4TNw2fMyz')

UPLOAD_FOLDER = 'uploads'
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)

def extract_text_from_pdf(pdf_path):
    text = ""
    with pdfplumber.open(pdf_path) as pdf:
        for page in pdf.pages:
            text += page.extract_text()
    return text

def extract_text_from_docx(docx_path):
    doc = docx.Document(docx_path)
    text = []
    for para in doc.paragraphs:
        text.append(para.text)
    return "\n".join(text)

def pptx_to_text(pptx_file):
    prs = Presentation(pptx_file)
    all_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                all_text.append(shape.text)
    return "\n".join(all_text)

def transcribe_audio(api_key, audio_file):
    return "Transcribed audio text here..."

def convert_to_text(file, file_type):
    if file_type == 'pdf':
        return extract_text_from_pdf(file)
    elif file_type == 'docx':
        return extract_text_from_docx(file)
    elif file_type == 'pptx':
        return pptx_to_text(file)
    elif file_type == 'audio':
        api_key = "your-audio-api-key"  
        return transcribe_audio(api_key, file)
    else:
        return None

@app.route("/api/generate_flashcards", methods=['POST'])
def generate_flashcards():
    if 'file' in request.files:
        file = request.files['file']
        file_type = file.filename.rsplit('.', 1)[1].lower()  
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], file.filename)
        file.save(file_path)
        text = convert_to_text(file_path, file_type)
    else:
        data = request.get_json()
        text = data.get('text', '')

    if not text:
        return jsonify({'error': 'Text is missing or invalid file type'}), 400

    try:
        summary_response = co.generate(
            model='command-xlarge-nightly',
            prompt=f"Summarize the following text:\n\n{text}",
            max_tokens=150,
            temperature=0.5,
        )
        summary = summary_response.generations[0].text.strip()

        questions_response = co.generate(
            model='command-xlarge-nightly',
            prompt=f"Create a list of questions from the following summary:\n\n{summary}",
            max_tokens=1000,
            temperature=0.7,
        )
        questions = [gen.text.strip() for gen in questions_response.generations]

        flashcards = []
        for question in questions:
            answer_response = co.generate(
                model='command-xlarge-nightly',
                prompt=f"Provide a concise answer to the following question based on the text: {text}\n\nQuestion: {question}",
                max_tokens=1000,
                temperature=0.7,
            )
            answer = answer_response.generations[0].text.strip()
            flashcards.append({
                "question": question,
                "answer": answer
            })

        return jsonify({'flashcards': flashcards}), 200

    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route("/api/generate_quiz", methods=['POST'])
def generate_quiz():
    if 'file' in request.files:
        file = request.files['file']
        file_type = file.filename.rsplit('.', 1)[1].lower()  
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], file.filename)
        file.save(file_path)
        text = convert_to_text(file_path, file_type)
    else:
        data = request.get_json()
        text = data.get('text', '')

    if not text:
        return jsonify({'error': 'Text is missing or invalid file type'}), 400

    try:
        # Generate quiz questions
        questions_response = co.generate(
            model='command-xlarge-nightly',
            prompt=f"Generate multiple-choice quiz questions from the following text:\n\n{text}",
            max_tokens=1000,
            temperature=0.7,
        )
        questions = [gen.text.strip() for gen in questions_response.generations]

        quiz = []
        for question in questions:
            # Generate the correct answer
            correct_answer_response = co.generate(
                model='command-xlarge-nightly',
                prompt=f"Provide the correct answer for the following quiz question:\n\nQuestion: {question}",
                max_tokens=1000,
                temperature=0.7,
            )
            correct_answer = correct_answer_response.generations[0].text.strip()

            # Generate incorrect answers
            incorrect_answers_response = co.generate(
                model='command-xlarge-nightly',
                prompt=f"Provide 3 plausible but incorrect answers for the following quiz question:\n\nQuestion: {question}",
                max_tokens=1000,
                temperature=0.7,
            )
            incorrect_answers = [gen.text.strip() for gen in incorrect_answers_response.generations[:3]]

            # Combine the correct answer with the incorrect answers
            answers = incorrect_answers + [correct_answer]

            # Shuffle the answers to randomize their order
            random.shuffle(answers)

            quiz.append({
                "question": question,
                "correct_answer": correct_answer,
                "answers": answers  # Now the answers array contains the correct and incorrect answers
            })

        return jsonify({'quiz': quiz}), 200


    except Exception as e:
        return jsonify({'error': str(e)}), 500

if __name__ == "__main__":
    app.run(debug=True, port=8080)