from dotenv import load_dotenv
from pprint import pprint
from docx import Document
from pptx import Presentation
import cohere
import random
import requests
import genanki
import pdfplumber
import assemblyai as aai
import os

load_dotenv()


def get_cohere_token():
    return os.getenv("COHERE_API_KEY")


def get_assembly_ai_token():
    return os.getenv("ASSEMBLY_AI_API_KEY")


def feed_cohere_text(text):

    co = cohere.Client(get_cohere_token())

    summary_generated = co.generate(
        model="command-xlarge-nightly",
        prompt=f"Summarize the following text:\n\n{text}",
        max_tokens=150,
        num_generations=1,
        temperature=0.5,
    )
    summary = summary_generated.generations[0].text.strip()
    questions_response = co.generate(
        model="command-xlarge-nightly",
        prompt=f"Create a list of questions from the following summary. Each question should start with -.\n\n{summary}",
        max_tokens=100,
        num_generations=5,
        temperature=0.7,
    )
    questions = [gen.text.strip() for gen in questions_response.generations]
    reranked_questions = co.rerank(
        query=f"Select the most relevant questions from the following list:",
        documents=questions,
        top_n=3,
        return_documents=True,
    )

    best_result = max(reranked_questions.results, key=lambda x: x.relevance_score)
    text_with_questions = best_result.document.text
    selected_questions = [
        line.strip()
        for line in text_with_questions.splitlines()
        if line.startswith("-")
    ]
    # for question in questions:
    #     print(question)
    q_a_dictionary = []
    for question in selected_questions:
        answer_response = co.generate(
            model="command-xlarge-nightly",
            prompt=f"Provide a concise answer to the following question based on the text: {text}\n\nQuestion: {question}",
            max_tokens=100,
            num_generations=1,
            temperature=0.7,
        )
        answer = answer_response.generations[0].text.strip()
        q_a_dictionary.append({"question": question, "answer": answer})

    formatted_output = ""
    for item in q_a_dictionary:
        formatted_output += f"{item['question']}\n{item['answer']}\n\n"
    return formatted_output


def generate_apkg_from_text(text, deck_name="unnamed"):
    model = genanki.Model(
        random.randrange(1 << 30, 1 << 31),
        "Simple Model",
        fields=[
            {"name": "Question"},
            {"name": "Answer"},
        ],
        templates=[
            {
                "name": "Card 1",
                "qfmt": "{{Question}}",
                "afmt": '{{FrontSide}}<hr id="answer">{{Answer}}',
            },
        ],
    )

    # Create a new deck
    deck = genanki.Deck(
        random.randrange(1 << 30, 1 << 31),
        deck_name,
    )

    # Split the text into lines
    lines = text.splitlines()

    # Initialize variables for questions and answers
    question = None
    answer = None

    # Iterate over lines and create note (card) pairs
    for line in lines:
        line = line.strip()
        if not line:
            # Ignore empty lines
            continue

        if question is None:
            question = line  # First line is the question
        elif answer is None:
            answer = line  # Second line is the answer
            # Add a new note (card) to the deck
            note = genanki.Note(model=model, fields=[question, answer])
            deck.add_note(note)
            # Reset for next question-answer pair
            question = None
            answer = None

    # Create the package and save it as a .apkg file
    genanki.Package(deck).write_to_file(f"{deck_name}.apkg")


# A function to get the text from an uplaoded PDF file
def extract_text_from_pdf(pdf_path):
    text = ""
    with pdfplumber.open(pdf_path) as pdf:
        for page in pdf.pages:
            text += page.extract_text()
    return text


def transcribe_audio(audio_file):

    aai.settings.api_key = get_assembly_ai_token()
    transcriber = aai.Transcriber()
    transcript = transcriber.transcribe(audio_file)
    return transcript.text


def extract_text_from_docx(file_path):
    doc = Document(file_path)
    full_text = []
    for para in doc.paragraphs:
        full_text.append(para.text)
    return "\n".join(full_text)


def extract_text_from_pptx(file_path):
    presentation = Presentation(file_path)
    full_text = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                full_text.append(shape.text)
    return "\n".join(full_text)


if __name__ == "__main__":
    # CLI for testing text input
    print("\n*** Anki  ***\n")

    deck_name = input("Enter the name of the deck you wish to create: ")
    text_to_process = input("\nEnter the content to generate questions from: ")

    # Call the Cohere API to generate accurate questions and answers
    questions_and_answers = feed_cohere_text(text_to_process)

    # Generate the Anki Deck from our formatted text
    generate_apkg_from_text(questions_and_answers, deck_name=deck_name)

    print("\n")
    print(f"Anki deck has been created as '{deck_name}.apkg'.")
    print("\n")
    # pprint(questions_and_answers)
